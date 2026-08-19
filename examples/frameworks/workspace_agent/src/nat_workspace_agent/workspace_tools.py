# SPDX-FileCopyrightText: Copyright (c) 2026, NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0

"""Workspace file tools for benchmarks that hand an agent a directory and a brief."""

import logging
import os
import subprocess
import sys
from collections.abc import AsyncGenerator
from pathlib import Path

from pydantic import Field

from nat.builder.builder import Builder
from nat.builder.function_info import FunctionInfo
from nat.cli.register_workflow import register_function
from nat.data_models.function import FunctionBaseConfig

MAX_READ_CHARS = 20000
CENSUS_ROWS = 40


def _root() -> Path:
    return Path(os.environ.get("NAT_WORKSPACE_DIR", ".")).resolve()


def _resolve(rel: str) -> Path:
    # A brief names the workspace by absolute path, so tasks echo it back whole, in part, or not
    # at all; each form folds under the root, which a task must never read or write outside of.
    root = _root()
    parts = Path(rel.strip()).parts
    parts = parts[parts.index(root.name) + 1:] if root.name in parts else tuple(
        x for x in parts if x not in ("/", "", "."))
    # Containment is judged lexically: resolve() would follow a link the workspace itself placed
    # (worlds are laid out as links into shared read-only data) and read it as an escape.
    p = Path(os.path.normpath(root.joinpath(*parts)))
    if p != root and root not in p.parents:
        raise ValueError(f"path escapes workspace: {rel}")
    return p


logger = logging.getLogger(__name__)

# Extraction runs out of process: one malformed PDF in a real workspace can hang pdfminer for
# minutes or crash the interpreter, and neither is catchable where the tool call happens.
_PDF_CHILD = ("import sys, pdfplumber\n"
              "with pdfplumber.open(sys.argv[1]) as pdf:\n"
              "    sys.stdout.write('\\n'.join((p.extract_text() or '') for p in pdf.pages[:40]))\n")


def _pdf_text(p: Path) -> str | None:
    try:
        done = subprocess.run([sys.executable, "-c", _PDF_CHILD, str(p)],
                              capture_output=True, timeout=60, check=False)
    except subprocess.TimeoutExpired:
        logger.warning("PDF %s took over 60s to parse and was skipped", p.name)
        return None
    if done.returncode != 0:
        logger.warning("PDF %s could not be parsed (exit %s)", p.name, done.returncode)
        return None
    return done.stdout.decode("utf-8", "ignore").strip() or None


def _extract(p: Path) -> str | None:
    """Text from a workspace file; None when the bytes carry no readable text.

    Office formats are zipped XML, so their text comes out with the stdlib alone.
    Reading them as UTF-8 instead yields mojibake that floods the context window.
    """
    import re
    import zipfile

    suffix = p.suffix.lower()
    if suffix in {".docx", ".xlsx", ".pptx"}:
        try:
            with zipfile.ZipFile(p) as z:
                parts = [n for n in z.namelist() if n.endswith(".xml") and "rels" not in n]
                chunks = []
                for name in parts[:40]:
                    raw = z.read(name).decode("utf-8", errors="ignore")
                    chunks.append(re.sub(r"<[^>]+>", " ", raw))
                return re.sub(r"\s+", " ", " ".join(chunks)).strip() or None
        except Exception:  # noqa: BLE001
            return None
    if suffix == ".pdf":
        return _pdf_text(p)
    # A binary whose first bytes happen to be ASCII still reads as mojibake, so the
    # extension decides rather than a null-byte probe alone.
    if suffix in {".doc", ".xls", ".ppt", ".png", ".jpg", ".jpeg", ".gif", ".zip", ".bin", ".so"}:
        return None
    try:
        data = p.read_bytes()
    except Exception:  # noqa: BLE001
        return None
    if b"\x00" in data[:4096]:
        return None
    return data.decode("utf-8", errors="replace")


class WorkspaceListConfig(FunctionBaseConfig, name="workspace_list"):
    max_entries: int = Field(default=200, description="Cap on returned paths")


@register_function(config_type=WorkspaceListConfig)
async def workspace_list(config: WorkspaceListConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """List files in the workspace."""

    async def _run(subdir: str = "", contains: str = "") -> str:
        base = _resolve(subdir) if subdir else _root()
        if not base.is_dir():
            return f"not a directory: {subdir}"
        needle = (contains or "").strip().lower()
        hits: list[str] = []
        census: dict[str, int] = {}
        for p in sorted(base.rglob("*")):
            if not p.is_file():
                continue
            rel = str(p.relative_to(_root()))
            if needle and needle not in rel.lower():
                continue
            census[str(Path(rel).parent)] = census.get(str(Path(rel).parent), 0) + 1
            hits.append(f"{rel}  ({p.stat().st_size} bytes)")
        if not hits:
            return f"(no file matches {contains!r})" if needle else "(empty)"
        if len(hits) <= config.max_entries:
            return "\n".join(hits)
        # An arbitrary slice of a large tree tells the agent nothing about where to look, so
        # past the cap the listing becomes a folder census short enough to survive any output cap.
        rows = sorted(census.items(), key=lambda kv: -kv[1])[:CENSUS_ROWS]
        return (f"{len(hits)} files match -- too many to list. Folders, largest first; "
                f"open one with `subdir`, or filter with `contains`:\n" +
                "\n".join(f"{d}/  ({n} files)" for d, n in rows))

    yield FunctionInfo.from_fn(
        _run,
        description=("List workspace files. Args: `subdir` relative to the root, and `contains` "
                     "to keep only paths holding that substring -- use it, the tree is large."))


class WorkspaceReadConfig(FunctionBaseConfig, name="workspace_read"):
    max_chars: int = Field(default=MAX_READ_CHARS, description="Cap on returned characters")


@register_function(config_type=WorkspaceReadConfig)
async def workspace_read(config: WorkspaceReadConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Read one workspace file as text."""

    async def _run(path: str, offset: int = 0) -> str:
        p = _resolve(path)
        if not p.is_file():
            return f"no such file: {path}"
        text = _extract(p)
        if text is None:
            return f"{path} is a binary file ({p.stat().st_size} bytes) with no text extractor."
        chunk = text[offset:offset + config.max_chars]
        rest = len(text) - offset - len(chunk)
        if rest <= 0:
            return chunk
        # Say where to continue: without it the model re-reads the same head forever, and a repeat
        # breaker cannot tell that apart from a loop. A differing offset is also a differing call.
        return f"{chunk}\n... {rest} more characters, call again with offset={offset + len(chunk)}"

    yield FunctionInfo.from_fn(_run, description=(
        "Read a workspace file as text. Args: `path` relative to the root, and `offset` to continue "
        "a long file from where the last call stopped."))


def _add_table(doc, rows: list[str]) -> None:
    cells = [[c.strip() for c in r.strip("|").split("|")] for r in rows]
    # The `|---|:--:|` rule under a markdown header carries no data, so it must not become a row.
    cells = [r for r in cells if not all(set(c) <= set("-: ") for c in r)]
    if not cells:
        return
    table = doc.add_table(rows=len(cells), cols=max(len(r) for r in cells))
    table.style = "Table Grid"
    for i, row in enumerate(cells):
        for j, value in enumerate(row):
            table.cell(i, j).text = value


def _write_docx(p: Path, text: str) -> str:
    """Build a real Word document out of markdown-ish text."""
    import re

    from docx import Document

    doc = Document()
    para: list[str] = []
    rows: list[str] = []

    def flush_para() -> None:
        if para:
            doc.add_paragraph(" ".join(para))
            para.clear()

    def flush_rows() -> None:
        if rows:
            _add_table(doc, rows)
            rows.clear()

    for raw in text.splitlines():
        line = raw.strip()
        if line.startswith("|") and line.endswith("|"):
            flush_para()
            rows.append(line)
            continue
        flush_rows()
        head = re.match(r"(#{1,3})\s+(.+)", line)
        bullet = re.match(r"([-*+\u2022]|\d+[.)])\s+", line)
        if head or bullet or not line:
            flush_para()
        if head:
            doc.add_heading(head[2].strip(" #"), level=len(head[1]))
        elif bullet:
            doc.add_paragraph(line)
        elif line:
            para.append(line)
    flush_rows()
    flush_para()
    doc.save(p)
    return f"{len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables"


def _write_xlsx(p: Path, text: str) -> str:
    """Build a real Excel workbook out of CSV text."""
    import csv
    import io

    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    # StringIO rather than splitlines(): only a real stream keeps a newline inside a quoted field.
    for row in csv.reader(io.StringIO(text)):
        ws.append(row)
    wb.save(p)
    return f"{ws.max_row} rows x {ws.max_column} columns"


def _write_pptx(p: Path, text: str) -> str:
    """Build a real PowerPoint deck, one slide per blank-line-separated block."""
    import re

    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    for block in re.split(r"\n[ \t]*\n", text.strip()):
        lines = [ln.strip().lstrip("#-*\u2022 ").strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = lines[0]
        slide.placeholders[1].text = "\n".join(lines[1:])
    prs.save(p)
    return f"{len(prs.slides)} slides"


class WorkspaceWriteConfig(FunctionBaseConfig, name="workspace_write"):
    pass


@register_function(config_type=WorkspaceWriteConfig)
async def workspace_write(config: WorkspaceWriteConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Write a deliverable into the workspace."""

    async def _run(path: str, content: str) -> str:
        p = _resolve(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        # Writing "to" a symlink was never meant to edit what it points at; a workspace laid out as
        # links into a shared read-only world would otherwise reject the write or corrupt the world.
        if p.is_symlink():
            p.unlink()
        build = {".docx": _write_docx, ".xlsx": _write_xlsx, ".pptx": _write_pptx}.get(p.suffix.lower())
        if build is None:
            p.write_text(content, encoding="utf-8")
            return f"wrote {p.relative_to(_root())} ({len(content)} chars)"
        try:
            detail = build(p, content)
        except Exception as exc:  # noqa: BLE001
            # Saving the text under the Office name instead would look like a delivered file and
            # score zero at grading time, so the agent has to hear about the failure now.
            return f"failed to build {p.name}: {exc}. Resend `content` in the shape that extension expects."
        return f"wrote {p.relative_to(_root())} ({detail})"

    yield FunctionInfo.from_fn(_run, description=(
        "Write a deliverable into the workspace. Args: `path` relative to the root, and `content`. "
        "The extension picks the format that gets built, so send `content` in the shape it expects: "
        "`.xlsx` wants CSV text -- header row first, one line per row, fields holding a comma quoted; "
        "`.docx` wants markdown-ish prose -- a blank line ends a paragraph, a leading `#`, `##` or `###` "
        "makes a heading of that level, and a run of `| a | b |` rows becomes a real table; "
        "`.pptx` wants one blank-line-separated block per slide, first line the title and the rest bullets. "
        "Every other extension is stored as the exact text you send."))


class WorkspaceSearchConfig(FunctionBaseConfig, name="workspace_search"):
    max_hits: int = Field(default=60, description="Cap on returned matching lines")


@register_function(config_type=WorkspaceSearchConfig)
async def workspace_search(config: WorkspaceSearchConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Find which workspace files hold a string, without reading each one whole."""

    async def _run(query: str, subdir: str = "", path_contains: str = "") -> str:
        needle = (query or "").strip().lower()
        if not needle:
            return "give a non-empty query"
        base = _resolve(subdir) if subdir else _root()
        if not base.is_dir():
            return f"not a directory: {subdir}"
        hits: list[str] = []
        for p in sorted(base.rglob("*")):
            if not p.is_file() or (path_contains and path_contains.lower() not in str(p).lower()):
                continue
            text = _extract(p)
            if text is None:
                continue
            rel = str(p.relative_to(_root()))
            for i, line in enumerate(text.splitlines(), 1):
                if needle in line.lower():
                    hits.append(f"{rel}:{i}: {line.strip()[:200]}")
                    if len(hits) > config.max_hits:
                        # Past the cap the agent needs a narrower query, not an arbitrary prefix.
                        return ("\n".join(hits[:config.max_hits]) +
                                f"\n... more than {config.max_hits} matches; narrow the query.")
        return "\n".join(hits) if hits else f"(no line contains {query!r})"

    yield FunctionInfo.from_fn(_run, description=(
        "Search workspace file contents for a string and return matching lines with their paths. "
        "Args: `query`, optional `subdir`, and `path_contains` to restrict which files are scanned."))
