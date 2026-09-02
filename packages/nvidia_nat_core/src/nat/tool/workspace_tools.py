# SPDX-FileCopyrightText: Copyright (c) 2026, NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0

"""Workspace file tools for benchmarks that hand an agent a directory and a brief."""

import contextlib
import logging
import os
import re
import shlex
import tempfile
import subprocess
import time
import sys
from collections.abc import AsyncGenerator
from collections import OrderedDict
from pathlib import Path

from pydantic import Field

from nat.builder.builder import Builder
from nat.builder.function_info import FunctionInfo
from nat.cli.register_workflow import register_function
from nat.data_models.function import FunctionBaseConfig

MAX_READ_CHARS = 20000
CENSUS_ROWS = 40


# Where a bridged session starts; the harness runs every container task set with this as its cwd.
_CONTAINER_ROOT = "/app"


def _root() -> Path:
    # No workspace was chosen, so scratch space -- the process cwd would hand over the whole checkout.
    if not os.environ.get("NAT_WORKSPACE_DIR"):
        os.environ["NAT_WORKSPACE_DIR"] = tempfile.mkdtemp(prefix="nat_workspace_")
    return Path(os.environ["NAT_WORKSPACE_DIR"]).resolve()


def _bridge() -> str | None:
    """The task's own container, when the harness opened one. The same flag run_code reads."""
    return os.environ.get("NAT_BRIDGE_URL") if os.environ.get("NAT_BRIDGE_READY") else None


def _sh(command: str, timeout: float = 60.0) -> tuple[bool, str]:
    """One shell command where the agent's files are. -> (ran, output).

    Shell and not python: six of terminalbench's fourteen local images carry no python at all, and
    the bridged session routes a command line to bash either way. `cd` first so it routes there.
    """
    import json as _json
    import urllib.request

    uri = _bridge()
    if not uri:
        return False, ""
    body = _json.dumps({"generated_code": f"cd {_CONTAINER_ROOT} && " + command,
                        "language": "python", "timeout": timeout}).encode()
    try:
        with urllib.request.urlopen(
                urllib.request.Request(uri.rstrip("/") + "/execute", body,
                                       {"Content-Type": "application/json"}),
                timeout=timeout + 15) as answer:
            got = _json.loads(answer.read())
    except Exception as exc:  # noqa: BLE001 -- unreachable and refusing mean the same thing here
        return True, f"the workspace is unreachable right now ({type(exc).__name__})."
    return True, (got.get("stdout") or "") + (got.get("stderr") or "")


def _where() -> str:
    """Container briefs name their files absolutely, and the root IS that directory: `x` and
    `/app/x` are one file. Unsaid, the model reads the mismatch as "these tools cannot reach it"
    and writes through bash heredocs instead."""
    return (f" This task's workspace root is {_CONTAINER_ROOT}, so `x` and {_CONTAINER_ROOT}/x name "
            f"the same file; either form works here." if _bridge() else "")


def _bare(step: str) -> str:
    """A step as text: agents send them bulleted, numbered, or already boxed."""
    return re.sub(r"^[-*\d.)\s]*(\[[ xX-]\])?\s*", "", step.strip())


def _steps_in(text) -> list[str]:
    """The steps an agent sent, whether as a list or a line per step.

    `steps` is `list[str] = []`, and the missing `| None` is the whole point. A tool call arrives
    as XML, where every parameter is text; the qwen3_xml parser turns the text back into a list by
    reading `properties["steps"]["type"]`, and an optional list is `anyOf: [array, null]`, which
    has no `type` at that level. The parser then falls back to string and hands the JSON array
    through as its own source text. That is how this went wrong the first time: splitlines() read
    the array as one step, the agent could not find its own plan in what came back, and one run
    rewrote it 245 times without ever writing the file it was asked for. Measured against the
    server: optional gives a string three times in three, plain `list[str]` a list three in three.
    `done` and `giving_up` are genuinely strings, and take the line-per-item path.
    """
    if isinstance(text, list):
        return [str(x).strip() for x in text if str(x).strip()]
    return [line for line in (text or "").splitlines() if line.strip()]


def _nothing_read(path: str, size: int, offset: int) -> str:
    """An empty read means an empty file or an offset past the end, and the bare "" it returned
    means those and a broken tool alike -- one run asked for the same empty file twice over."""
    if size <= 0:
        return f"{path} exists and is empty; there is nothing in it to read."
    return (f"nothing at offset {offset}: {path} is {size} characters long, so that is past its "
            f"end. Read it from 0, or from an offset below {size}.")


def _resolve(rel: str) -> Path:
    # Briefs echo the root path whole, partial or not at all; each form folds under the root.
    root = _root()
    parts = Path(rel.strip()).parts
    parts = parts[parts.index(root.name) + 1:] if root.name in parts else tuple(
        x for x in parts if x not in ("/", "", "."))
    # Lexical containment: resolve() would follow the links the workspace itself placed.
    p = Path(os.path.normpath(root.joinpath(*parts)))
    if p != root and root not in p.parents:
        raise ValueError(f"path escapes workspace: {rel}")
    return p


logger = logging.getLogger(__name__)

_PART_OF = ("\n\n[only the first {kept} {unit} were read, of {whole}; the rest is not shown and searching this file will not find it]")

# Counted where the cap middleware counts its own: a dropped match is a dropped match.
from nat.middleware.output_limit.output_limit_middleware import FIRED

# Out of process: one malformed PDF can hang pdfminer or crash the interpreter, uncatchably.
_PDF_CHILD = ("import sys, pdfplumber\n"
              "with pdfplumber.open(sys.argv[1]) as pdf:\n"
              "    sys.stdout.write('\\n'.join((p.extract_text() or '') for p in pdf.pages[:40]))\n"
              # Said, not guessed: without the count the parent cannot tell 40 pages from 400,
              # and workspace_read then reports the head of a long document as the whole of it.
              "    sys.stderr.write(str(len(pdf.pages)))\n")


# A PDF that could not be parsed cannot be parsed the next time either, and each attempt costs
# the full timeout. Two workspace runs spent their wall clock re-parsing the same broken file.
_PDF_REFUSED: dict[tuple[str, int, float], None] = {}
# Successful extractions too: a second search over the same tree re-parsed 1950 PDFs from scratch.
_EXTRACTED: "OrderedDict[tuple[str, int, float], str | None]" = OrderedDict()
_EXTRACT_CACHE_MAX = 4000


def _pdf_text(p: Path) -> str | None:
    try:
        stamp = (str(p), p.stat().st_size, p.stat().st_mtime)
    except OSError:
        stamp = (str(p), -1, -1.0)
    if stamp in _PDF_REFUSED:
        return None
    # A workspace tree carries empty placeholder files; an empty one is not a broken PDF, and
    # warning per file per search buried the reads that did fail.
    if stamp[1] == 0:
        return None
    try:
        done = subprocess.run([sys.executable, "-c", _PDF_CHILD, str(p)],
                              capture_output=True, timeout=60, check=False)
    except subprocess.TimeoutExpired:
        logger.warning("PDF %s took over 60s to parse and was skipped", p.name)
        _PDF_REFUSED[stamp] = None
        return None
    if done.returncode != 0:
        logger.warning("PDF %s could not be parsed (exit %s)", p.name, done.returncode)
        _PDF_REFUSED[stamp] = None
        return None
    text = done.stdout.decode("utf-8", "ignore").strip() or None
    # The child wrote the page census to stderr: 40 pages of a 266-page filing is not the filing,
    # and workspace_read reports a head it was never told was a head as the whole document.
    if text:
        with contextlib.suppress(ValueError):
            whole = int(done.stderr.decode("utf-8", "ignore").strip() or 0)
            if whole > 40:
                text += _PART_OF.format(kept=40, whole=whole, unit="pages of this PDF")
    return text


def _extract(p: Path) -> str | None:
    """Text from a workspace file, or None: Office formats are zipped XML read via the stdlib;
    reading them as UTF-8 yields mojibake that floods the context window."""
    try:
        stamp = (str(p), p.stat().st_size, p.stat().st_mtime)
    except OSError:
        stamp = None
    if stamp is not None and stamp in _EXTRACTED:
        _EXTRACTED.move_to_end(stamp)
        return _EXTRACTED[stamp]
    text = _extract_uncached(p)
    if stamp is not None:
        _EXTRACTED[stamp] = text
        while len(_EXTRACTED) > _EXTRACT_CACHE_MAX:
            _EXTRACTED.popitem(last=False)
    return text


def _extract_uncached(p: Path) -> str | None:
    import re
    import zipfile

    suffix = p.suffix.lower()
    if suffix in {".docx", ".xlsx", ".pptx"}:
        try:
            with zipfile.ZipFile(p) as z:
                parts = [n for n in z.namelist() if n.endswith(".xml") and "rels" not in n]
                chunks = []
                for name in parts[:40]:
                    raw = z.read(name).decode("utf-8", errors="ignore")
                    chunks.append(re.sub(r"<[^>]+>", " ", raw))
                text = re.sub(r"\s+", " ", " ".join(chunks)).strip() or None
                # A document read down to its first 40 parts is not the document, and silence here
                # reads downstream as "this is all of it".
                if text and len(parts) > 40:
                    text += _PART_OF.format(kept=40, whole=len(parts), unit="parts of this file")
                return text
        except Exception:
            return None
    if suffix == ".pdf":
        text = _pdf_text(p)
        # An unparseable .pdf that is really plain text is an agent-written deliverable: let it read back.
        if text is not None:
            return text
    # ASCII-leading binaries still read as mojibake, so the extension decides, not a byte probe.
    if suffix in {".doc", ".xls", ".ppt", ".png", ".jpg", ".jpeg", ".gif", ".zip", ".bin", ".so"}:
        return None
    try:
        data = p.read_bytes()
    except Exception:
        return None
    if b"\x00" in data[:4096]:
        return None
    return data.decode("utf-8", errors="replace")


class WorkspaceListConfig(FunctionBaseConfig, name="workspace_list"):
    max_entries: int = Field(default=200, description="Cap on returned paths")


def _built(target: Path, content: str) -> str:
    """Build `content` into `target` in the shape its extension asks for. -> what was built."""
    build = {".docx": _write_docx, ".xlsx": _write_xlsx, ".pptx": _write_pptx}.get(target.suffix.lower())
    if build is None:
        target.write_text(content, encoding="utf-8")
        return f"{content.count(chr(10)) + 1} lines"
    return build(target, content)


def _put(path: str, blob: bytes) -> tuple[bool, str]:
    """Place bytes at `path` in the container. Built here and carried over as base64, because the
    xlsx, docx and pptx writers need libraries the task's own image has no reason to hold."""
    import base64
    q = shlex.quote(path)
    encoded = base64.b64encode(blob).decode()
    return _sh(f"mkdir -p \"$(dirname {q})\" && printf %s {shlex.quote(encoded)} | base64 -d > {q} "
               f"&& echo __WROTE__")


def _formatted(root, pairs, contains: str, max_entries: int) -> str:
    """One listing format for both trees, the local one and the container's: the files, or past the
    cap a folder census, which is the only thing an arbitrary slice of them could not tell."""
    needle = (contains or "").strip().lower()
    hits: list[str] = []
    census: dict[str, int] = {}
    for rel, size in pairs:
        if not rel or (needle and needle not in rel.lower()):
            continue
        census[str(Path(rel).parent)] = census.get(str(Path(rel).parent), 0) + 1
        hits.append(f"{rel}  ({size} bytes)" if size else rel)
    # The absolute root, so code run in a sandbox can open these files by path.
    head = f"workspace root: {root}"
    if not hits:
        return head + ("\n(no file matches %r)" % contains if needle else "\n(empty)")
    if len(hits) <= max_entries:
        return "\n".join([head, *hits])
    rows = sorted(census.items(), key=lambda kv: -kv[1])[:CENSUS_ROWS]
    return (f"{head}\n{len(hits)} files match -- too many to list. Folders, largest first; "
            f"open one with `subdir`, or filter with `contains`:\n" +
            "\n".join(f"{d}/  ({n} files)" for d, n in rows))


def _from_find(out: str, where: str, contains: str, max_entries: int) -> str:
    """`find -printf '%P\\t%s'` output, formatted the way a local tree is."""
    pairs = []
    for line in out.splitlines():
        name, _, size = line.partition("\t")
        pairs.append((f"{where}/{name}".lstrip("./") if where not in (".", "") else name, size))
    return _formatted(_CONTAINER_ROOT, pairs, contains, max_entries)


def _listing(subdir: str = "", contains: str = "", max_entries: int = 200) -> str:
    """Module level so a test can reach it: the census branch shipped broken with nothing covering it."""
    base = _resolve(subdir) if subdir else _root()
    if not base.is_dir():
        return f"not a directory: {subdir}"
    pairs = [(str(p.relative_to(_root())), p.stat().st_size)
             for p in sorted(base.rglob("*")) if p.is_file()]
    return _formatted(_root(), pairs, contains, max_entries)


@register_function(config_type=WorkspaceListConfig)
async def workspace_list(config: WorkspaceListConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """List files in the workspace."""

    async def _run(subdir: str = "", contains: str = "") -> str:
        if _bridge():
            where = subdir.strip("/ ") or "."
            ran, out = _sh(f"find {shlex.quote(where)} -type f -printf '%P\\t%s\\n' 2>/dev/null "
                           f"|| find {shlex.quote(where)} -type f")
            if ran:
                return _from_find(out, where, contains, config.max_entries)
        return _listing(subdir, contains, config.max_entries)

    yield FunctionInfo.from_fn(
        _run,
        description=("List workspace files. Args: `subdir` relative to the root, and `contains` "
                     "to keep only paths holding that substring -- use it, the tree is large."))


class WorkspaceReadConfig(FunctionBaseConfig, name="workspace_read"):
    max_chars: int = Field(default=MAX_READ_CHARS, description="Cap on returned characters")


@register_function(config_type=WorkspaceReadConfig)
async def workspace_read(config: WorkspaceReadConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Read one workspace file as text."""

    async def _run(path: str, offset: int = 0) -> str:
        if _bridge():
            q = shlex.quote(path)
            ran, out = _sh(f"if [ -d {q} ]; then echo __DIR__; elif [ -f {q} ]; then "
                           f"wc -c < {q}; echo __CUT__; tail -c +{offset + 1} {q} | head -c "
                           f"{config.max_chars}; else echo __MISSING__; fi")
            if ran:
                if out.startswith("__DIR__"):
                    return f"{path} is a directory -- list it with workspace_list, or name a file in it."
                if out.startswith("__MISSING__"):
                    return f"no such file: {path}"
                total, _, chunk = out.partition("__CUT__\n")
                try:
                    size = int(total.strip())
                except ValueError:
                    return out
                if not chunk:
                    return _nothing_read(path, size, offset)
                rest = size - offset - len(chunk)
                if rest <= 0:
                    return chunk
                return (f"{chunk}\n... {rest} more characters, call again with "
                        f"offset={offset + len(chunk)}")
        p = _resolve(path)
        if p.is_dir():
            return f"{path} is a directory -- list it with workspace_list, or name a file in it."
        if not p.is_file():
            return f"no such file: {path}"
        text = _extract(p)
        if text is None:
            return f"{path} is a binary file ({p.stat().st_size} bytes) with no text extractor."
        chunk = text[offset:offset + config.max_chars]
        if not chunk:
            return _nothing_read(path, len(text), offset)
        rest = len(text) - offset - len(chunk)
        if rest <= 0:
            return chunk
        # Say where to continue, or the model re-reads the same head and looks like a repeat loop.
        return f"{chunk}\n... {rest} more characters, call again with offset={offset + len(chunk)}"

    yield FunctionInfo.from_fn(_run, description=(
        "Read a workspace file as text. Args: `path` relative to the root, and `offset` to continue "
        "a long file from where the last call stopped." + _where()))


def _add_table(doc, rows: list[str]) -> None:
    cells = [[c.strip() for c in r.strip("|").split("|")] for r in rows]
    # The `|---|:--:|` rule under a markdown header carries no data, so it must not become a row.
    cells = [r for r in cells if not all(set(c) <= set("-: ") for c in r)]
    if not cells:
        return
    table = doc.add_table(rows=len(cells), cols=max(len(r) for r in cells))
    table.style = "Table Grid"
    for i, row in enumerate(cells):
        for j, value in enumerate(row):
            table.cell(i, j).text = value


def _write_docx(p: Path, text: str) -> str:
    """Build a real Word document out of markdown-ish text."""
    import re

    from docx import Document

    doc = Document()
    para: list[str] = []
    rows: list[str] = []

    def flush_para() -> None:
        if para:
            doc.add_paragraph(" ".join(para))
            para.clear()

    def flush_rows() -> None:
        if rows:
            _add_table(doc, rows)
            rows.clear()

    for raw in text.splitlines():
        line = raw.strip()
        if line.startswith("|") and line.endswith("|"):
            flush_para()
            rows.append(line)
            continue
        flush_rows()
        head = re.match(r"(#{1,3})\s+(.+)", line)
        bullet = re.match(r"([-*+\u2022]|\d+[.)])\s+", line)
        if head or bullet or not line:
            flush_para()
        if head:
            doc.add_heading(head[2].strip(" #"), level=len(head[1]))
        elif bullet:
            doc.add_paragraph(line)
        elif line:
            para.append(line)
    flush_rows()
    flush_para()
    doc.save(p)
    return f"{len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables"


def _write_xlsx(p: Path, text: str) -> str:
    """Build a real Excel workbook out of CSV text."""
    import csv
    import io
    import re

    from openpyxl import Workbook

    def _typed(cell: str):
        """A number written as text is one Excel flags and every formula skips, and a column holding
        both sorts worse than one holding either. Plain decimals convert; a leading zero, a plus or
        an exponent means an identifier -- 007, +1, 1e5 -- and stays the string it was sent as."""
        body = cell.strip()
        if not re.fullmatch(r"-?(?:0|[1-9][0-9]*)(?:\.[0-9]+)?", body):
            return cell
        return float(body) if "." in body else int(body)

    wb = Workbook()
    ws = wb.active
    # StringIO rather than splitlines(): only a real stream keeps a newline inside a quoted field.
    for row in csv.reader(io.StringIO(text)):
        ws.append([_typed(c) for c in row])
    wb.save(p)
    return f"{ws.max_row} rows x {ws.max_column} columns"


def _write_pptx(p: Path, text: str) -> str:
    """Build a real PowerPoint deck, one slide per blank-line-separated block."""
    import re

    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    for block in re.split(r"\n[ \t]*\n", text.strip()):
        lines = [ln.strip().lstrip("#-*\u2022 ").strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = lines[0]
        slide.placeholders[1].text = "\n".join(lines[1:])
    prs.save(p)
    return f"{len(prs.slides)} slides"


class WorkspaceWriteConfig(FunctionBaseConfig, name="workspace_write"):
    pass


@register_function(config_type=WorkspaceWriteConfig)
async def workspace_write(config: WorkspaceWriteConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Write a deliverable into the workspace."""

    async def _run(path: str, content: str) -> str:
        if _bridge():
            import tempfile as _tmp
            with _tmp.TemporaryDirectory(prefix="ws-out-") as staging:
                local = Path(staging) / Path(path).name
                try:
                    detail = _built(local, content)
                except Exception as exc:  # noqa: BLE001 -- the builder's own words beat a traceback
                    return f"could not build {path}: {exc}"
                ran, out = _put(path, local.read_bytes())
            if ran:
                return (f"wrote {path} ({detail})" if "__WROTE__" in out
                        else out.strip() or f"could not write {path}")
        p = _resolve(path)
        # An empty or directory `path` writes onto the directory itself, and that OS error carries
        # an absolute host path the caller can do nothing with.
        if p == _root() or p.is_dir():
            return f"`path` must name a file inside the workspace; {path!r} names a directory."
        p.parent.mkdir(parents=True, exist_ok=True)
        # Writing to a symlink never means editing its target; linked worlds would reject or corrupt.
        if p.is_symlink():
            p.unlink()
        build = {".docx": _write_docx, ".xlsx": _write_xlsx, ".pptx": _write_pptx}.get(p.suffix.lower())
        if build is None:
            # Said at the moment the lines go: a description telling the agent to prefer
            # workspace_edit moved 1 framework in 4, and the loss is silent otherwise.
            before = p.read_text(encoding="utf-8", errors="ignore").count("\n") + 1 if p.is_file() else 0
            try:
                p.write_text(content, encoding="utf-8")
            except OSError as exc:
                return f"could not write {path}: {exc.strerror or exc}."
            after = content.count("\n") + 1
            note = (f" It held {before} lines and now holds {after}; the other {before - after} are "
                    "gone. If that was not intended, workspace_edit changes one passage and leaves "
                    "the rest." if before > after + max(5, before // 10) else "")
            return f"wrote {p.relative_to(_root())} ({len(content)} chars){note}"
        try:
            detail = build(p, content)
        except Exception as exc:
            # Text saved under an Office name looks delivered and grades zero; the failure must be heard.
            return f"failed to build {p.name}: {exc}. Resend `content` in the shape that extension expects."
        return f"wrote {p.relative_to(_root())} ({detail})"

    yield FunctionInfo.from_fn(_run, description=(
        "Write a deliverable into the workspace, replacing whatever was there. To change part of a "
        "file that already exists, use workspace_edit instead -- everything you do not resend here "
        "is lost. Args: `path` relative to the root, and `content`. "
        "The extension picks the format that gets built, so send `content` in the shape it expects: "
        "`.xlsx` wants CSV text -- header row first, one line per row, fields holding a comma quoted; "
        "`.docx` wants markdown-ish prose -- a blank line ends a paragraph, a leading `#`, `##` or `###` "
        "makes a heading of that level, and a run of `| a | b |` rows becomes a real table; "
        "`.pptx` wants one blank-line-separated block per slide, first line the title and the rest bullets. "
        "Every other extension is stored as the exact text you send." + _where()))


class WorkspaceSearchConfig(FunctionBaseConfig, name="workspace_search"):
    max_hits: int = Field(default=60, description="Cap on returned matching lines")
    # max_hits bounds the answer, not the work: a query that matches nothing still opened every
    # PDF in the tree. One workspace holds 1950 of them and the scan measured 81 minutes.
    max_seconds: float = Field(default=90.0, description="Wall clock one search may spend scanning")
    max_documents: int = Field(default=250, description="Documents whose text may be extracted per search")


@register_function(config_type=WorkspaceSearchConfig)
async def workspace_search(config: WorkspaceSearchConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Find which workspace files hold a string, without reading each one whole."""

    async def _run(query: str, subdir: str = "", path_contains: str = "") -> str:
        needle = (query or "").strip().lower()
        if not needle:
            return "give a non-empty query"
        if _bridge():
            where = shlex.quote(subdir.strip("/ ") or ".")
            keep = f" | grep -F -- {shlex.quote(path_contains)}" if path_contains else ""
            ran, out = _sh(f"grep -rniI -F -- {shlex.quote(query)} {where} 2>/dev/null{keep} "
                           f"| head -n {config.max_hits + 1}")
            if ran:
                rows = [r for r in out.splitlines() if r.strip()]
                if not rows:
                    return f"(no line contains {query!r})"
                if len(rows) > config.max_hits:
                    return ("\n".join(rows[:config.max_hits])
                            + f"\n... more than {config.max_hits} matches; narrow the query.")
                return "\n".join(rows)
        base = _resolve(subdir) if subdir else _root()
        if not base.is_dir():
            return f"not a directory: {subdir}"
        hits: list[str] = []
        started, opened, seen, skipped = time.time(), 0, 0, 0
        for p in sorted(base.rglob("*")):
            if not p.is_file() or (path_contains and path_contains.lower() not in str(p).lower()):
                continue
            seen += 1
            costly = p.suffix.lower() in {".pdf", ".docx", ".xlsx", ".pptx"}
            if costly and (opened >= config.max_documents
                           or time.time() - started > config.max_seconds):
                skipped += 1
                continue
            opened += costly
            text = _extract(p)
            if text is None:
                continue
            rel = str(p.relative_to(_root()))
            for i, line in enumerate(text.splitlines(), 1):
                if needle in line.lower():
                    said = line.strip()
                    # A CSV row or a minified file is one long line, and its first 200 characters
                    # read exactly like all of it. The bridged grep does not cut at all, so
                    # unmarked this tool answers differently depending on whether a container is up.
                    hits.append(f"{rel}:{i}: {said[:200]}"
                                + (f"  ...[+{len(said) - 200} more on this line]"
                                   if len(said) > 200 else ""))
                    if len(hits) > config.max_hits:
                        # Past the cap the agent needs a narrower query, not an arbitrary prefix.
                        FIRED[config.type] += 1
                        return ("\n".join(hits[:config.max_hits]) +
                                f"\n... more than {config.max_hits} matches; narrow the query.")
        # Said, not hidden: a truncated scan that reads as "no matches" sends the agent away
        # from the file it was looking for.
        note = (f"\n[scanned {seen - skipped} of {seen} files in {time.time() - started:.0f}s; "
                f"{skipped} documents were left unopened -- narrow with subdir= or path_contains=]"
                if skipped else "")
        return ("\n".join(hits) + note) if hits else (f"(no line contains {query!r})" + note)

    yield FunctionInfo.from_fn(_run, description=(
        "Search workspace file contents and return matching lines with their paths. `query` is plain "
        "text matched case-insensitively, not a regular expression -- for a regex use bash with grep. "
        "Args: `query`, optional `subdir`, and `path_contains` to restrict which files are scanned." + _where()))


class WorkspaceEditConfig(FunctionBaseConfig, name="workspace_edit"):
    pass


@register_function(config_type=WorkspaceEditConfig)
async def workspace_edit(config: WorkspaceEditConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Replace one exact passage in a file, rather than rewriting the file around it."""

    async def _run(path: str, old: str, new: str) -> str:
        if _bridge():
            q = shlex.quote(path)
            ran, body = _sh(f"if [ -f {q} ]; then cat {q}; else echo __MISSING__; fi")
            if ran:
                if body.startswith("__MISSING__"):
                    return f"{path} is not a file in the workspace; workspace_list shows what is."
                seen = body.count(old)
                if seen == 0:
                    return f"that exact text is not in {path}; read it again and copy the passage."
                if seen > 1:
                    return f"that text appears {seen} times in {path}; include more of it."
                done, out = _put(path, body.replace(old, new, 1).encode())
                if done and "__WROTE__" in out:
                    return f"edited {path} ({len(old)} chars -> {len(new)})"
                return out.strip() or f"could not write {path}"
        p = _resolve(path)
        if not p.is_file():
            return f"{path} is not a file in the workspace; workspace_list shows what is."
        try:
            body = p.read_text(encoding="utf-8")
        except (OSError, UnicodeDecodeError) as exc:
            return f"could not read {path}: {getattr(exc, 'strerror', None) or exc}."
        hits = body.count(old)
        if hits == 0:
            return (f"that passage does not appear in {path}; read it first and copy the text "
                    "exactly, whitespace included.")
        if hits > 1:
            # Editing the first of several is how a file quietly gets the wrong one changed.
            return f"that passage appears {hits} times in {path}; extend `old` until it is unique."
        # As workspace_write does: a staged world is symlinks into the shared dataset, and writing
        # through one edits the dataset for every run after this.
        if p.is_symlink():
            p.unlink()
        p.write_text(body.replace(old, new), encoding="utf-8")
        return f"edited {p.relative_to(_root())} ({len(old)} chars -> {len(new)})"

    yield FunctionInfo.from_fn(_run, description=(
        "Replace one exact passage inside a file, leaving the rest untouched. Use this to change an "
        "existing file; workspace_write replaces the whole file and loses anything you did not "
        "resend.\n\nArgs:\n    path (str): the file, relative to the workspace root.\n"
        "    old (str): the exact text to replace, unique in the file.\n"
        "    new (str): what to put there." + _where()))


class WorkspaceShellConfig(FunctionBaseConfig, name="workspace_shell"):
    uri: str = Field(default="http://127.0.0.1:6000", description="Sandbox base URL")
    timeout: float = Field(default=60.0, description="Seconds one command may run")
    max_output_characters: int = Field(default=16000, description="Truncate combined output here")


@register_function(config_type=WorkspaceShellConfig)
async def workspace_shell(config: WorkspaceShellConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """A shell in the sandbox, rooted at the workspace."""
    import json as _json

    import httpx

    async def _run(command: str) -> str:
        # A bridged sandbox is the task's own container: its session already starts where the work
        # is and routes a command line to bash itself, so the host path baked in below would name a
        # directory that does not exist there. This is what run_code's own preamble check does.
        if os.environ.get("NAT_BRIDGE_READY"):
            wrapper = command
        else:
            # Run through the sandbox, never on this host: the tool exists because agents were
            # wrapping subprocess in run_code to get here anyway, and that path had no root and no cap.
            # The cut is caught here, not left to surface as a TimeoutExpired traceback: that reads
            # as a crash, and it throws away what the command had already printed.
            wrapper = (
                "import subprocess, os\n"
                f"cwd = {_root().as_posix()!r}\n"
                "os.makedirs(cwd, exist_ok=True)\n"
                "try:\n"
                f"    r = subprocess.run({command!r}, shell=True, cwd=cwd, capture_output=True,\n"
                f"                       text=True, timeout={config.timeout})\n"
                "except subprocess.TimeoutExpired as cut:\n"
                # TimeoutExpired carries the output as bytes even under text=True, and stderr as None.
                "    got = [cut.stdout, cut.stderr]\n"
                "    print(''.join(p.decode(errors='replace') if isinstance(p, bytes) else (p or '')\n"
                "                  for p in got), end='')\n"
                f"    print('\\n[stopped at the {config.timeout:g}s limit -- anything after this "
                "was not run]', end='')\n"
                "else:\n"
                "    print(r.stdout, end='')\n"
                "    print(r.stderr, end='')\n"
                "    print(f'\\n[exit {r.returncode}]' if r.returncode else '', end='')\n")
        try:
            async with httpx.AsyncClient(timeout=config.timeout + 15) as client:
                answer = await client.post(config.uri.rstrip("/") + "/execute",
                                           json={"generated_code": wrapper,
                                                 "timeout": config.timeout, "language": "python"})
                answer.raise_for_status()
                body = answer.json()
        except Exception as exc:  # noqa: BLE001 -- any failure to run means the same thing
            # Not "right now": that reads as a wait, and an agent told to wait re-sends the same
            # command until its budget is gone -- thirty-six times in one run measured here. What
            # it cannot work out for itself is that nothing it types will change this answer.
            return (f"the shell is not running ({type(exc).__name__}). This is the sandbox, not "
                    f"your command: the same command will get this same answer, so use the other "
                    f"tools and say so in your answer if the task needed a shell.")
        out = (body.get("stdout") or "") + (body.get("stderr") or "")
        if body.get("process_status") not in (None, "completed", "success"):
            out = f"[{body['process_status']}]\n{out}"
        # A command that succeeds and prints nothing -- a heredoc, a bare `#` comment the model
        # meant as a thought -- otherwise returns "", which reads as a tool that did nothing and
        # gets re-sent. One run re-sent the same comment 95 times. Say the command ran and left
        # no output, so re-running it is pointless; think() is where a thought belongs.
        out = out.strip() or ("the command ran and exited 0 with no output -- re-running it will "
                              "return this again. If you meant to reason, use think; otherwise move on.")
        cap = config.max_output_characters
        if len(out) <= cap:
            return out
        FIRED[config.type] += 1
        return out[:cap] + (f"\n...[cut at {cap} characters. Send the output to a file and read it, "
                            "or filter it here with head, tail or grep -- running this again returns "
                            "the same cut.]")

    # Named, not just described: agents guessed /app, the image's own workdir, and lost 3 steps.
    yield FunctionInfo.from_fn(_run, description=(
        "Run one shell command in the workspace and return its output. Each call is a new shell, "
        "so a `cd` or an exported variable is gone by the next one -- chain them in one command "
        "line instead. The working directory is "
        f"the workspace root, {_root().as_posix()}, so paths are relative to it. For anything on "
        "the web use web_search and web_fetch rather than curl or urllib here: those keep what "
        "they read where the rest of the run can see it.\n\n"
        "Args:\n    command (str): the command line, e.g. `ls -la` or `python -m pytest -q`."))


class ThinkConfig(FunctionBaseConfig, name="think"):
    pass


@register_function(config_type=ThinkConfig)
async def think(config: ThinkConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Somewhere to reason mid-task. Measured by Anthropic at +54% relative on tau-bench airline."""

    async def _run(thought: str) -> str:
        return "Noted. Continue."

    yield FunctionInfo.from_fn(_run, description=(
        "Think a step through without acting. Nothing happens and nothing is fetched or changed; "
        "use it to work out what a tool just told you, check a rule before you act on it, or plan "
        "the next few steps.\n\nArgs:\n    thought (str): the reasoning to work through."))


# Keyed by workspace and shared by every agent on one task. A process running tasks in turn reuses
# one workspace, so the key alone does not keep them apart -- reset_task_plan drops the plan at each
# task boundary rather than carrying one task's steps into the next.
_PLANS: dict[str, list[str]] = {}


def reset_task_plan(workspace: str = "") -> None:
    """Drop the task_list plan for a workspace, at a task boundary."""
    _PLANS.pop(str(Path(workspace).resolve()) if workspace else str(_root()), None)


class TaskListConfig(FunctionBaseConfig, name="task_list"):
    pass


@register_function(config_type=TaskListConfig)
async def task_list(config: TaskListConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """The plan as run state, shared by every agent on this task.

    State, not a workspace file: one there was searched, graded, and half-visible over the bridge.

    Three states, not two: closing a step by SOLVING it is the only way an agent had, so a step it
    could not solve stayed open and it kept trying. `giving_up` closes one and says why, and the
    reason is what stops the next agent -- or the next turn of this one -- repeating the attempt.
    """

    async def _run(steps: list[str] = [], done: str = "", giving_up: str = "",
                   because: str = "") -> str:
        key = str(_root())
        lines = list(_PLANS.get(key, []))
        before = list(lines)
        asked = _steps_in(steps)
        if asked:
            # A closed step stays closed: each specialist rewrites this list, and a plain replace
            # reopened what the one before it had already finished.
            shut = {l[6:].split("  (")[0].strip().lower(): l for l in lines if not l.startswith("- [ ]")}
            # Steps often arrive already bulleted, and "- [ ] - step" reads as a broken list.
            lines = [shut.get(_bare(s).strip().lower(), f"- [ ] {_bare(s)}")
                     for s in asked]
        missed = []
        for mark, box, why in ([(x, "- [x]", "") for x in _steps_in(done)]
                               + [(x, "- [-]", because) for x in _steps_in(giving_up)]):
            for i, line in enumerate(lines):
                if line.startswith("- [ ]") and mark.strip().lower() in line.lower():
                    lines[i] = line.replace("- [ ]", box, 1) + (f"  ({why.strip()})" if why.strip() else "")
                    break
            else:
                missed.append(mark.strip())
        if lines:
            _PLANS[key] = lines
        if not lines:
            return "The list is empty; send `steps` to start one."
        left = sum(l.startswith("- [ ]") for l in lines)
        gave = sum(l.startswith("- [-]") for l in lines)
        tail = f"({left} of {len(lines)} still open" + (f", {gave} given up on)" if gave else ")")
        # Saying nothing changed is the whole point: a silent no-op reads as success and gets
        # retried -- one run sent the same plan 19 times and read back the same words every time.
        if lines == before and (asked or done.strip() or giving_up.strip()):
            tail += "\nNothing changed: this is the list as it already stood."
        if missed:
            tail += ("\nNo open step matches " + ", ".join(repr(m) for m in missed)
                     + " -- already closed, or never on the list.")
        return "\n".join(lines) + "\n\n" + tail

    yield FunctionInfo.from_fn(_run, description=(
        "Keep the plan for this task as a checklist. Call it once with `steps` to write the plan, "
        "then with `done` after finishing one, or with `giving_up` and `because` for one you tried "
        "and could not settle -- writing that down is what keeps you, and anyone after you, from "
        "trying it again. With no arguments it shows where you are. Every agent here shares it."
        "\n\nArgs:\n    steps (list[str]): the plan, one step per item -- replaces any existing "
        "list.\n"
        "    done (str): text identifying a step to tick off, one per line.\n"
        "    giving_up (str): text identifying a step to close unsolved.\n"
        "    because (str): why that step could not be settled."))


class FinishConfig(FunctionBaseConfig, name="finish"):
    pass


@register_function(config_type=FinishConfig)
async def finish(config: FinishConfig, builder: Builder) -> AsyncGenerator[FunctionInfo, None]:
    """Stopping, as an action. Without one an agent keeps picking the best tool it has."""
    from nat.middleware.agent_finish import AgentFinished

    async def _run(answer: str) -> str:
        raise AgentFinished(answer)

    yield FunctionInfo.from_fn(_run, description=(
        "Finish, giving your answer. Call this when the task is done -- checking work you have "
        "already checked cannot change it.\n\n"
        "Args:\n    answer (str): the complete answer, in the form the task asked for."))
